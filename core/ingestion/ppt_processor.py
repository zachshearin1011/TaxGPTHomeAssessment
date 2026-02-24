from __future__ import annotations

import logging
import struct
import subprocess
import tempfile
from pathlib import Path

from core.ingestion.csv_processor import Document, GraphTriple
from core.ingestion.pdf_processor import TAX_CONCEPTS

logger = logging.getLogger(__name__)


class PPTProcessor:

    def __init__(self, path: Path, chunk_size: int = 512, chunk_overlap: int = 64):
        self.path = path
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def extract_text(self) -> list[tuple[int, str]]:
        suffix = self.path.suffix.lower()
        if suffix == ".pptx":
            return self._extract_pptx()
        elif suffix == ".ppt":
            return self._extract_ppt_legacy()
        else:
            raise ValueError(f"Unsupported format: {suffix}")

    def to_documents(self) -> list[Document]:
        slides = self.extract_text()
        docs: list[Document] = []

        for slide_num, slide_text in slides:
            chunks = self._chunk_text(slide_text)
            for i, chunk in enumerate(chunks):
                if len(chunk.strip()) < 20:
                    continue
                docs.append(Document(
                    text=chunk,
                    metadata={
                        "source_file": self.path.name,
                        "slide": slide_num,
                        "chunk_index": i,
                    },
                    source=self.path.name,
                    doc_type="ppt",
                ))

        logger.info("Produced %d document chunks from %s", len(docs), self.path.name)
        return docs

    def to_graph_triples(self) -> list[GraphTriple]:
        slides = self.extract_text()
        full_text = " ".join(t for _, t in slides).lower()

        triples: list[GraphTriple] = []

        found: list[str] = []
        for concept in TAX_CONCEPTS:
            if concept in full_text:
                found.append(concept)
                triples.append(GraphTriple(
                    subject=self.path.stem,
                    subject_type="Presentation",
                    predicate="covers",
                    object=concept,
                    object_type="TaxConcept",
                ))

        for i, c1 in enumerate(found):
            for c2 in found[i + 1:]:
                triples.append(GraphTriple(
                    subject=c1, subject_type="TaxConcept",
                    predicate="related_to",
                    object=c2, object_type="TaxConcept",
                ))

        logger.info("Extracted %d graph triples from %s", len(triples), self.path.name)
        return triples

    def _extract_pptx(self) -> list[tuple[int, str]]:
        from pptx import Presentation

        prs = Presentation(str(self.path))
        slides: list[tuple[int, str]] = []

        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            parts.append(row_text)
            if parts:
                slides.append((i, "\n".join(parts)))

        logger.info("Extracted text from %d slides (pptx) of %s", len(slides), self.path.name)
        return slides

    def _extract_ppt_legacy(self) -> list[tuple[int, str]]:
        converted = self._try_libreoffice_convert()
        if converted:
            return converted
        return self._extract_ppt_binary()

    def _try_libreoffice_convert(self) -> list[tuple[int, str]] | None:
        for cmd in ("libreoffice", "soffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
            try:
                with tempfile.TemporaryDirectory() as tmp:
                    result = subprocess.run(
                        [cmd, "--headless", "--convert-to", "pptx", "--outdir", tmp, str(self.path)],
                        capture_output=True, timeout=120,
                    )
                    if result.returncode == 0:
                        pptx_path = Path(tmp) / (self.path.stem + ".pptx")
                        if pptx_path.exists():
                            logger.info("Converted %s to pptx via LibreOffice", self.path.name)
                            from pptx import Presentation
                            prs = Presentation(str(pptx_path))
                            slides: list[tuple[int, str]] = []
                            for i, slide in enumerate(prs.slides, 1):
                                parts: list[str] = []
                                for shape in slide.shapes:
                                    if shape.has_text_frame:
                                        for para in shape.text_frame.paragraphs:
                                            text = para.text.strip()
                                            if text:
                                                parts.append(text)
                                if parts:
                                    slides.append((i, "\n".join(parts)))
                            return slides
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
        return None

    def _extract_ppt_binary(self) -> list[tuple[int, str]]:
        import olefile

        if not olefile.isOleFile(str(self.path)):
            logger.warning("Not a valid OLE file: %s", self.path.name)
            return [(1, "")]

        ole = olefile.OleFileIO(str(self.path))
        texts: list[str] = []

        if ole.exists("PowerPoint Document"):
            stream = ole.openstream("PowerPoint Document").read()
            texts = self._parse_ppt_records(stream)
        elif ole.exists("Current User"):
            for stream_name in ole.listdir():
                joined = "/".join(stream_name)
                if "powerpoint" in joined.lower():
                    data = ole.openstream(stream_name).read()
                    texts.extend(self._parse_ppt_records(data))

        ole.close()

        if not texts:
            logger.warning("No text extracted from binary PPT: %s", self.path.name)
            return [(1, "")]

        slide_text = "\n".join(texts)
        logger.info(
            "Extracted %d text fragments from binary PPT %s",
            len(texts),
            self.path.name,
        )
        return [(1, slide_text)]

    @staticmethod
    def _parse_ppt_records(data: bytes) -> list[str]:
        texts: list[str] = []

        def _walk(offset: int, end: int) -> None:
            pos = offset
            while pos < end - 8:
                rec_ver_inst = struct.unpack_from("<H", data, pos)[0]
                rec_ver = rec_ver_inst & 0x0F
                rec_type = struct.unpack_from("<H", data, pos + 2)[0]
                rec_len = struct.unpack_from("<I", data, pos + 4)[0]

                if rec_len > end - pos - 8:
                    break

                if rec_ver == 0x0F:
                    _walk(pos + 8, pos + 8 + rec_len)
                else:
                    if rec_type == 0x0FA0 and rec_len > 0:
                        try:
                            text = data[pos + 8: pos + 8 + rec_len].decode("utf-16-le", errors="replace")
                            cleaned = text.strip().replace("\r", "\n").replace("\x00", "")
                            if cleaned and len(cleaned) > 1:
                                texts.append(cleaned)
                        except Exception:
                            pass
                    elif rec_type == 0x0FA8 and rec_len > 0:
                        try:
                            text = data[pos + 8: pos + 8 + rec_len].decode("latin-1", errors="replace")
                            cleaned = text.strip().replace("\r", "\n")
                            if cleaned and len(cleaned) > 1:
                                texts.append(cleaned)
                        except Exception:
                            pass

                pos += 8 + rec_len

        _walk(0, len(data))
        return texts

    def _chunk_text(self, text: str) -> list[str]:
        words = text.split()
        chunks: list[str] = []
        start = 0
        while start < len(words):
            end = start + self.chunk_size
            chunk = " ".join(words[start:end])
            chunks.append(chunk)
            start += self.chunk_size - self.chunk_overlap
        return chunks